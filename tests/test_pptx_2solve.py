"""Gerador PPTX padrão 2Solve: estrutura, conteúdo e formatação BRL."""

from pptx import Presentation

from src.connectors.pptx_2solve import (
    ProposalContent,
    ProposalItem,
    brl,
    build_proposal_pptx,
    proposal_filename,
)


def conteudo_exemplo() -> ProposalContent:
    return ProposalContent(
        cliente="Usina X",
        projeto="Adequação de Instrumentação",
        titulo="Modernização da malha de medição fiscal",
        subtitulo="Conformidade metrológica e visibilidade em tempo real",
        numero="2001999",
        revisao="A",
        escopo=[
            "Levantar e revisar a lista de instrumentos da área 300",
            "Elaborar P&ID as-built das malhas FT-301 e FT-302",
        ],
        itens=[
            ProposalItem(descricao="2S Tools Gateway LoRaWAN", quantidade=1,
                         valor_unitario=12512.74),
            ProposalItem(descricao="Transmissor IoT MD Universal", quantidade=40,
                         valor_unitario=1094.59),
            ProposalItem(descricao="Comissionamento em campo", quantidade=1,
                         unidade="VB", valor_unitario=None),
        ],
        prazo_entrega="90 dias após pedido de compra",
        condicoes_pagamento=["28 dias após faturamento", "Frete FOB"],
        observacoes=["Impostos inclusos conforme regime de Lucro Presumido."],
        vendedor="Thiago Pompermayer",
    )


def todos_os_textos(path) -> str:
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def test_gera_pptx_com_todas_as_secoes(tmp_path):
    content = conteudo_exemplo()
    path = build_proposal_pptx(content, tmp_path / proposal_filename(content))

    assert path.exists()
    prs = Presentation(str(path))
    assert len(prs.slides) == 5  # capa, escopo, itens, condições, contato

    texto = todos_os_textos(path)
    assert "Modernização da malha de medição fiscal" in texto
    assert "USINA X" in texto  # rodapé/capa em caixa alta
    assert "FT-301" in texto
    assert "Transmissor IoT MD Universal" in texto
    assert "sob consulta" in texto  # item sem valor não vira preço inventado
    assert "R$ 56.296,34" in texto  # total só dos itens cotados
    assert "90 dias após pedido de compra" in texto
    assert "28 dias após faturamento" in texto
    assert "10.821.258/0001-02" in texto  # dados ToSolve
    assert "comercial@2solve.com" in texto


def test_pptx_sem_escopo_nem_itens_pula_slides_opcionais(tmp_path):
    content = ProposalContent(
        cliente="Cliente Y", projeto="Projeto Z", titulo="Estudo preliminar"
    )
    path = build_proposal_pptx(content, tmp_path / "minima.pptx")
    assert len(Presentation(str(path)).slides) == 3  # capa, condições, contato


def test_brl_formata_padrao_brasileiro():
    assert brl(56296.34) == "R$ 56.296,34"
    assert brl(1000) == "R$ 1.000,00"


def test_nome_de_arquivo_sem_caracteres_problematicos():
    content = conteudo_exemplo()
    assert proposal_filename(content) == "Proposta_2Solve_Usina_X_2001999_RevA.pptx"
