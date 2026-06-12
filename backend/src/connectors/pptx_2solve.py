"""Gerador de proposta PPTX no padrão visual 2Solve.

Identidade extraída de templates/Petroreconcavo_Carnauba_rev00.pptx:
fundo escuro #0A0E1A, cards #1A2333, acento ciano #03E8DC, azul #007AFF,
textos #F4F4F4/#C6C6C6/#A8A8A8, fonte Work Sans, rodapé "CLIENTE • PROJETO".
Dados fixos da ToSolve conforme propostas oficiais (templates/*.pdf).
"""

from __future__ import annotations

import re
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pydantic import BaseModel, Field

# ----- identidade visual -----

BG = RGBColor(0x0A, 0x0E, 0x1A)
CARD = RGBColor(0x1A, 0x23, 0x33)
CYAN = RGBColor(0x03, 0xE8, 0xDC)
BLUE = RGBColor(0x00, 0x7A, 0xFF)
TEXT = RGBColor(0xF4, 0xF4, 0xF4)
MUTED = RGBColor(0xC6, 0xC6, 0xC6)
DIM = RGBColor(0xA8, 0xA8, 0xA8)
FONT = "Work Sans"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

LOGO_PATH = Path(__file__).resolve().parents[2] / "templates" / "assets" / "logo_2solve.png"

COMPANY = {
    "razao": "ToSolve Engenharia e Tecnologia Ltda",
    "site": "2solve.com",
    "cnpj": "10.821.258/0001-02",
    "endereco": "Av. Adalberto Simão Nader, 675 — República, Vitória/ES — CEP 29070-010",
    "telefone": "(27) 3026-3806",
    "email": "comercial@2solve.com",
}


# ----- conteúdo da proposta -----


class ProposalItem(BaseModel):
    descricao: str
    quantidade: float = 1
    unidade: str = "UN"
    valor_unitario: float | None = None  # None => "sob consulta"

    @property
    def valor_total(self) -> float | None:
        if self.valor_unitario is None:
            return None
        return self.quantidade * self.valor_unitario


class ProposalContent(BaseModel):
    cliente: str
    projeto: str
    titulo: str
    subtitulo: str = ""
    numero: str = ""
    revisao: str = "0"
    escopo: list[str] = Field(default_factory=list)
    itens: list[ProposalItem] = Field(default_factory=list)
    prazo_entrega: str = "A combinar"
    condicoes_pagamento: list[str] = Field(default_factory=list)
    validade_dias: int = 30
    observacoes: list[str] = Field(default_factory=list)
    vendedor: str = ""


def brl(value: float) -> str:
    formatted = f"{value:,.2f}".replace(",", "@").replace(".", ",").replace("@", ".")
    return f"R$ {formatted}"


def slugify(text: str) -> str:
    return re.sub(r"[^A-Za-z0-9]+", "_", text).strip("_")


def proposal_filename(content: ProposalContent) -> str:
    numero = content.numero or date.today().strftime("%Y%m%d")
    return f"Proposta_2Solve_{slugify(content.cliente)}_{numero}_Rev{content.revisao}.pptx"


# ----- helpers de desenho -----


def _fill_bg(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    size: float,
    color: RGBColor = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def _card(slide, left, top, width, height, color: RGBColor = CARD):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _header(slide, kicker: str, title: str, subtitle: str, page: int) -> None:
    """Padrão de cabeçalho do template: kicker ciano, título grande, subtítulo."""
    _card(slide, Inches(0.55), Inches(0.52), Inches(0.07), Inches(0.95), CYAN)
    _text(slide, Inches(0.75), Inches(0.45), Inches(9.5), Inches(0.35),
          kicker.upper(), 13, CYAN, bold=True)
    _text(slide, Inches(0.75), Inches(0.78), Inches(10.5), Inches(0.6),
          title, 24, TEXT, bold=True)
    if subtitle:
        _text(slide, Inches(0.75), Inches(1.32), Inches(10.5), Inches(0.4),
              subtitle, 13, MUTED)
    _text(slide, Inches(12.3), Inches(0.45), Inches(0.8), Inches(0.5),
          f"{page:02d}", 24, CYAN, bold=True)


def _footer(slide, content: ProposalContent) -> None:
    _text(slide, Inches(0.55), Inches(7.05), Inches(9.0), Inches(0.3),
          f"{content.cliente.upper()}  •  {content.projeto.upper()}", 9.5, DIM)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(11.9), Inches(6.9),
                                 height=Inches(0.45))


def _new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    _fill_bg(slide)
    return slide


# ----- slides -----


def _slide_cover(prs: Presentation, content: ProposalContent) -> None:
    slide = _new_slide(prs)
    _card(slide, Emu(0), Inches(6.9), SLIDE_W, Inches(0.08), CYAN)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Inches(0.55), Inches(0.5),
                                 height=Inches(0.8))
    _text(slide, Inches(0.55), Inches(2.3), Inches(11.0), Inches(0.4),
          "PROPOSTA TÉCNICA-COMERCIAL", 14, CYAN, bold=True)
    _text(slide, Inches(0.55), Inches(2.75), Inches(12.2), Inches(1.6),
          content.titulo, 42, TEXT, bold=True)
    if content.subtitulo:
        _text(slide, Inches(0.55), Inches(4.35), Inches(11.5), Inches(0.9),
              content.subtitulo, 17, MUTED)
    numero = f"Nº {content.numero}  •  " if content.numero else ""
    _text(slide, Inches(0.55), Inches(5.9), Inches(11.5), Inches(0.4),
          f"{content.cliente.upper()}  •  {content.projeto.upper()}", 16, CYAN, bold=True)
    _text(slide, Inches(0.55), Inches(6.35), Inches(11.5), Inches(0.35),
          f"{numero}Rev. {content.revisao}  •  {date.today().strftime('%d/%m/%Y')}",
          12, DIM)


def _slide_scope(prs: Presentation, content: ProposalContent, page: int) -> None:
    slide = _new_slide(prs)
    _header(slide, "Escopo de fornecimento", "O que está incluído",
            content.subtitulo, page)
    top = 2.0
    for bullet in content.escopo:
        _card(slide, Inches(0.75), Inches(top + 0.08), Inches(0.12), Inches(0.12), CYAN)
        _text(slide, Inches(1.05), Inches(top - 0.06), Inches(11.3), Inches(0.55),
              bullet, 14, MUTED)
        top += 0.62
    _footer(slide, content)


def _slide_items(prs: Presentation, content: ProposalContent, page: int) -> None:
    slide = _new_slide(prs)
    _header(slide, "Investimento", "Itens da proposta", "", page)

    headers = ["Item", "Descrição", "Qtd.", "Un.", "Valor unit.", "Valor total"]
    widths = [Inches(0.7), Inches(6.0), Inches(0.9), Inches(0.8), Inches(1.9), Inches(1.9)]
    table_shape = slide.shapes.add_table(
        len(content.itens) + 1, len(headers), Inches(0.55), Inches(1.95),
        Inches(12.2), Inches(0.4 * (len(content.itens) + 1)),
    )
    table = table_shape.table
    for col, width in enumerate(widths):
        table.columns[col].width = width

    def set_cell(row: int, col: int, value: str, *, bold=False,
                 color: RGBColor = MUTED, fill: RGBColor = CARD) -> None:
        cell = table.cell(row, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        cell.text = value
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = FONT
        paragraph.font.size = Pt(11)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color

    for col, header in enumerate(headers):
        set_cell(0, col, header, bold=True, color=TEXT, fill=RGBColor(0x1E, 0x29, 0x3B))
    total_geral = 0.0
    tem_sob_consulta = False
    for row, item in enumerate(content.itens, start=1):
        unit = brl(item.valor_unitario) if item.valor_unitario is not None else "sob consulta"
        if item.valor_total is None:
            total, tem_sob_consulta = "—", True
        else:
            total = brl(item.valor_total)
            total_geral += item.valor_total
        set_cell(row, 0, f"{row:02d}")
        set_cell(row, 1, item.descricao)
        set_cell(row, 2, f"{item.quantidade:g}")
        set_cell(row, 3, item.unidade)
        set_cell(row, 4, unit)
        set_cell(row, 5, total)

    label = "Investimento total" + (" (itens cotados)" if tem_sob_consulta else "")
    bottom = 1.95 + 0.42 * (len(content.itens) + 1) + 0.35
    _text(slide, Inches(7.0), Inches(bottom), Inches(3.3), Inches(0.4),
          label, 13, MUTED, align=PP_ALIGN.RIGHT)
    _text(slide, Inches(10.3), Inches(bottom - 0.08), Inches(2.45), Inches(0.5),
          brl(total_geral), 20, CYAN, bold=True, align=PP_ALIGN.RIGHT)
    _footer(slide, content)


def _slide_conditions(prs: Presentation, content: ProposalContent, page: int) -> None:
    slide = _new_slide(prs)
    _header(slide, "Condições comerciais", "Prazo, pagamento e validade", "", page)

    _card(slide, Inches(0.55), Inches(2.0), Inches(5.95), Inches(4.4))
    _text(slide, Inches(0.85), Inches(2.25), Inches(5.4), Inches(0.35),
          "PRAZO DE ENTREGA", 12, CYAN, bold=True)
    _text(slide, Inches(0.85), Inches(2.7), Inches(5.4), Inches(1.0),
          content.prazo_entrega, 14, TEXT)
    _text(slide, Inches(0.85), Inches(3.9), Inches(5.4), Inches(0.35),
          "VALIDADE DA PROPOSTA", 12, CYAN, bold=True)
    _text(slide, Inches(0.85), Inches(4.35), Inches(5.4), Inches(0.8),
          f"{content.validade_dias} ({_por_extenso(content.validade_dias)}) dias "
          "após a data de emissão.", 14, TEXT)

    _card(slide, Inches(6.8), Inches(2.0), Inches(5.95), Inches(4.4))
    _text(slide, Inches(7.1), Inches(2.25), Inches(5.4), Inches(0.35),
          "CONDIÇÕES DE PAGAMENTO", 12, CYAN, bold=True)
    top = 2.7
    for cond in content.condicoes_pagamento or ["A combinar."]:
        _text(slide, Inches(7.1), Inches(top), Inches(5.4), Inches(0.5),
              f"•  {cond}", 12.5, MUTED)
        top += 0.52
    _footer(slide, content)


def _slide_notes_contact(prs: Presentation, content: ProposalContent, page: int) -> None:
    slide = _new_slide(prs)
    _header(slide, "Informações finais", "Observações e contato", "", page)
    top = 2.0
    for obs in content.observacoes:
        _card(slide, Inches(0.75), Inches(top + 0.08), Inches(0.12), Inches(0.12), BLUE)
        _text(slide, Inches(1.05), Inches(top - 0.06), Inches(11.3), Inches(0.5),
              obs, 12.5, MUTED)
        top += 0.55

    card_top = max(top + 0.3, 4.4)
    _card(slide, Inches(0.55), Inches(card_top), Inches(12.2), Inches(1.9))
    _text(slide, Inches(0.85), Inches(card_top + 0.2), Inches(11.5), Inches(0.35),
          COMPANY["razao"].upper(), 13, CYAN, bold=True)
    contato = (
        f"CNPJ {COMPANY['cnpj']}  •  {COMPANY['endereco']}\n"
        f"{COMPANY['telefone']}  •  {COMPANY['email']}  •  {COMPANY['site']}"
    )
    if content.vendedor:
        contato += f"\nContato comercial: {content.vendedor}"
    _text(slide, Inches(0.85), Inches(card_top + 0.6), Inches(11.5), Inches(1.2),
          contato, 12, MUTED)
    _footer(slide, content)


def _por_extenso(n: int) -> str:
    nomes = {15: "quinze", 30: "trinta", 45: "quarenta e cinco", 60: "sessenta",
             90: "noventa"}
    return nomes.get(n, str(n))


# ----- API do connector -----


def build_proposal_pptx(content: ProposalContent, output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, content)
    page = 2
    if content.escopo:
        _slide_scope(prs, content, page)
        page += 1
    if content.itens:
        _slide_items(prs, content, page)
        page += 1
    _slide_conditions(prs, content, page)
    page += 1
    _slide_notes_contact(prs, content, page)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
